# ============================================================================
# Prométhée — Assistant IA desktop
# ============================================================================
# Auteur  : Pierre COUGET
# Licence : GNU Affero General Public License v3.0 (AGPL-3.0)
#           https://www.gnu.org/licenses/agpl-3.0.html
# Année   : 2026
# ----------------------------------------------------------------------------
# Ce fichier fait partie du projet Prométhée.
# Vous pouvez le redistribuer et/ou le modifier selon les termes de la
# licence AGPL-3.0 publiée par la Free Software Foundation.
# ============================================================================

"""
tools/export_template_tools.py — Génération de documents conformes aux modèles organisationnels
================================================================================================

Outils exposés (3) :

  Word / docx (2) :
    - export_docx_template   : génère un document Word en appliquant un gabarit
                               (.docx) fourni par l'organisation. Les styles,
                               en-têtes, pieds de page, logos et marges du modèle
                               sont préservés. Le contenu est injecté via les styles
                               nommés du gabarit (Heading 1/2/3, Normal, List Bullet…)
                               ou via les signets (bookmarks) de substitution.
    - list_docx_template_styles : liste les styles disponibles dans un gabarit .docx,
                               utile avant d'appeler export_docx_template pour choisir
                               les bons noms de styles.

  Présentation / pptx (1) :
    - export_pptx_template   : génère une présentation PowerPoint en appliquant un
                               gabarit (.pptx). Les layouts, le thème graphique, les
                               couleurs et la typographie de l'organisation sont
                               préservés. Les slides sont créés depuis les layouts
                               du gabarit (par nom ou par index).

Conventions communes
────────────────────
  - template_path : chemin absolu ou relatif au home vers le fichier gabarit
                    (.docx ou .pptx). Le gabarit n'est jamais modifié.
  - output_path   : chemin de destination du document produit.
                    Si omis, créé dans ~/Exports/Prométhée/.
  - Retour        : dict JSON {"path": "...", "size_bytes": N, "status": "ok"}
  - En cas d'erreur : {"error": "message", "status": "error"}

Bonnes pratiques pour les gabarits
────────────────────────────────────
  Gabarit Word (.docx) :
    - Les styles nommés dans le gabarit (ex : « Titre organisation », « Corps texte »)
      sont réutilisables directement via le champ "style" de chaque section.
    - Pour les zones de substitution fixes (en-tête de courrier, cadre de référence),
      placer des signets Word nommés dans le gabarit. L'outil remplacera leur contenu.
    - Les en-têtes, pieds de page et sections du gabarit sont hérités automatiquement.

  Gabarit PowerPoint (.pptx) :
    - Chaque layout du gabarit est accessible par son nom (champ "layout_name")
      ou son index (champ "layout_index").
    - Les placeholders (titres, corps, images) du layout sont peuplés dans l'ordre
      ou par leur idx.
    - Le thème (couleurs, polices, fond) est hérité intégralement du gabarit.

Prérequis :
    pip install python-docx python-pptx
"""

import copy
import json
import shutil
import tempfile
from pathlib import Path
from typing import Any

from core.tools_engine import tool, set_current_family, _TOOL_ICONS

set_current_family("export_template_tools", "Export depuis modèles", "📋")

_TOOL_ICONS.update({
    "export_docx_template":        "📄",
    "export_pptx_template":        "📊",
    "list_docx_template_styles":   "🔍",
})

# ── Répertoire de sortie par défaut ──────────────────────────────────────────

_DEFAULT_EXPORT_DIR = Path.home() / "Exports" / "Prométhée"


def _resolve_output(output_path: str, default_name: str) -> Path:
    if output_path:
        p = Path(output_path).expanduser()
        if not p.is_absolute():
            p = Path.home() / p
    else:
        _DEFAULT_EXPORT_DIR.mkdir(parents=True, exist_ok=True)
        p = _DEFAULT_EXPORT_DIR / default_name
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _resolve_template(template_path: str) -> Path:
    p = Path(template_path).expanduser()
    if not p.is_absolute():
        p = Path.home() / p
    return p


def _ok(path: Path, extra: dict | None = None) -> str:
    r = {"status": "ok", "path": str(path), "size_bytes": path.stat().st_size}
    if extra:
        r.update(extra)
    return json.dumps(r, ensure_ascii=False)


def _err(msg: str) -> str:
    return json.dumps({"status": "error", "error": msg}, ensure_ascii=False)


# ═════════════════════════════════════════════════════════════════════════════
# UTILITAIRE — LISTER LES STYLES D'UN GABARIT DOCX
# ═════════════════════════════════════════════════════════════════════════════

@tool(
    name="list_docx_template_styles",
    description=(
        "Liste tous les styles disponibles dans un gabarit Word (.docx). "
        "À appeler AVANT export_docx_template pour connaître les noms exacts "
        "des styles de l'organisation (titres, corps, tableaux, listes…). "
        "Retourne un objet JSON avec les styles regroupés par type : "
        "paragraph (paragraphes), character (caractères), table (tableaux), "
        "numbering (listes numérotées). "
        "Utile aussi pour vérifier les signets (bookmarks) disponibles dans le gabarit."
    ),
    parameters={
        "type": "object",
        "properties": {
            "template_path": {
                "type": "string",
                "description": (
                    "Chemin vers le gabarit .docx de l'organisation. "
                    "Ex : ~/Modèles/modele_note.docx ou /home/user/gabarits/lettre.docx"
                )
            }
        },
        "required": ["template_path"]
    }
)
def list_docx_template_styles(template_path: str) -> str:
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from lxml import etree

        tpl_path = _resolve_template(template_path)
        if not tpl_path.exists():
            return _err(f"list_docx_template_styles : gabarit introuvable : {tpl_path}")
        if tpl_path.suffix.lower() != ".docx":
            return _err("list_docx_template_styles : le fichier doit être un .docx")

        doc = Document(str(tpl_path))

        # Styles regroupés par type
        styles_by_type: dict[str, list[dict]] = {
            "paragraph": [],
            "character": [],
            "table":     [],
            "numbering": [],
        }
        for style in doc.styles:
            stype = str(style.type.name).lower() if style.type else "unknown"
            if stype in styles_by_type:
                styles_by_type[stype].append({
                    "name":     style.name,
                    "style_id": style.style_id,
                    "built_in": style.builtin,
                })

        # Signets présents dans le document
        bookmarks = []
        for bm in doc.element.findall(f".//{qn('w:bookmarkStart')}"):
            name_attr = bm.get(qn("w:name"))
            if name_attr and not name_attr.startswith("_"):
                bookmarks.append(name_attr)

        result = {
            "template": str(tpl_path),
            "styles":   styles_by_type,
            "bookmarks": bookmarks,
            "total_styles": sum(len(v) for v in styles_by_type.values()),
        }
        return json.dumps(result, ensure_ascii=False, indent=2)

    except Exception as e:
        return _err(f"list_docx_template_styles : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# WORD — EXPORT DEPUIS GABARIT DOCX
# ═════════════════════════════════════════════════════════════════════════════

def _apply_style_safe(paragraph, style_name: str, doc) -> bool:
    """Applique un style à un paragraphe ; retourne True si le style existe."""
    try:
        doc.styles[style_name]  # lève KeyError si absent
        paragraph.style = style_name
        return True
    except (KeyError, Exception):
        return False


def _replace_bookmark_text(doc, bookmark_name: str, new_text: str) -> bool:
    """
    Remplace le contenu d'un signet Word par new_text.
    Retourne True si le signet a été trouvé et modifié.
    """
    from docx.oxml.ns import qn
    from lxml import etree

    # Trouver le bookmarkStart avec ce nom
    for bm_start in doc.element.iter(qn("w:bookmarkStart")):
        if bm_start.get(qn("w:name")) == bookmark_name:
            bm_id = bm_start.get(qn("w:id"))
            # Trouver le paragraphe parent
            parent = bm_start.getparent()
            if parent is None:
                continue
            # Effacer les runs entre bookmarkStart et bookmarkEnd
            found_start = False
            for child in list(parent):
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if child == bm_start:
                    found_start = True
                    continue
                if found_start:
                    if tag == "bookmarkEnd" and child.get(qn("w:id")) == bm_id:
                        break
                    if tag == "r":
                        parent.remove(child)
            # Insérer un run avec le nouveau texte après bookmarkStart
            ns_w = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            r_elem = etree.SubElement(parent, f"{{{ns_w}}}r")
            t_elem = etree.SubElement(r_elem, f"{{{ns_w}}}t")
            t_elem.text = new_text
            t_elem.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
            # Repositionner après bookmarkStart
            bm_idx = list(parent).index(bm_start)
            parent.remove(r_elem)
            parent.insert(bm_idx + 1, r_elem)
            return True
    return False


@tool(
    name="export_docx_template",
    description=(
        "Génère un document Word (.docx) CONFORME AU MODÈLE DE L'ORGANISATION "
        "en appliquant un gabarit (.docx) fourni. "
        "Les en-têtes, pieds de page, logos, marges et styles graphiques du gabarit "
        "sont INTÉGRALEMENT préservés — seul le contenu est injecté. "
        "À utiliser dès qu'un modèle officiel (lettre, note, rapport, arrêté, fiche…) "
        "est disponible dans l'organisation. "
        "Préférer cet outil à export_docx quand un gabarit existe. "
        "\n\n"
        "FONCTIONNEMENT : "
        "1. Le gabarit est copié tel quel (styles + mise en page). "
        "2. Si clear_body=true, le corps du gabarit est vidé avant injection du contenu. "
        "3. Le contenu (sections) est ajouté en utilisant les styles nommés du gabarit. "
        "4. Les substitutions de signets (bookmarks) remplacent des zones fixes "
        "   (ex : en-tête de courrier, champ référence, date). "
        "\n\n"
        "CONSEIL : appeler list_docx_template_styles avant pour connaître les noms "
        "exacts des styles disponibles dans le gabarit. "
        "\n\n"
        "Structure JSON d'une section : "
        '{"heading": "Titre", "level": 1, '
        '"style": "Titre 1",  '
        '"paragraphs": ["Para 1.", "Para 2."], '
        '"bullets": ["item 1", "item 2"], '
        '"bullet_style": "List Bullet", '
        '"table": {"headers": ["Col A"], "rows": [["val1"]]}, '
        '"page_break": false}'
    ),
    parameters={
        "type": "object",
        "properties": {
            "template_path": {
                "type": "string",
                "description": (
                    "Chemin vers le gabarit .docx de l'organisation. "
                    "Ex : ~/Modèles/note_de_service.docx"
                )
            },
            "document": {
                "type": "object",
                "description": (
                    "Structure du contenu à injecter dans le gabarit. "
                    "Champs de premier niveau : "
                    "title (str, optionnel — titre principal), "
                    "bookmarks (objet clé→valeur pour substituer des signets du gabarit, "
                    "  ex : {\"date\": \"12 juin 2026\", \"ref\": \"RH-2026-042\"}), "
                    "sections (liste d'objets section). "
                    "Chaque section accepte : "
                    "heading (str), level (int 1-3), "
                    "style (str — nom exact du style du gabarit, ex 'Titre 1', 'Corps texte'), "
                    "paragraphs (liste de str), content (str), intro (str), "
                    "bullets (liste de str), bullet_style (str, style de liste du gabarit), "
                    "table (objet {headers, rows}), page_break (bool)."
                )
            },
            "clear_body": {
                "type": "boolean",
                "description": (
                    "Si true (défaut), vide le corps du gabarit avant d'injecter le contenu "
                    "— utile pour les gabarits avec contenu d'exemple. "
                    "Si false, le contenu est ajouté APRÈS le contenu existant du gabarit "
                    "— utile pour les gabarits avec en-tête de courrier fixe à conserver."
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination (ex: ~/Documents/note_2026.docx). Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["template_path", "document"]
    }
)
def export_docx_template(
    template_path: str,
    document: dict,
    clear_body: bool = True,
    output_path: str = "",
    filename: str = ""
) -> str:
    try:
        from docx import Document
        from docx.shared import Pt
        from docx.oxml.ns import qn

        tpl_path = _resolve_template(template_path)
        if not tpl_path.exists():
            return _err(f"export_docx_template : gabarit introuvable : {tpl_path}")
        if tpl_path.suffix.lower() != ".docx":
            return _err("export_docx_template : le gabarit doit être un fichier .docx")

        name = filename or (document.get("title", "document") + ".docx")
        if not name.endswith(".docx"):
            name += ".docx"
        dest = _resolve_output(output_path, name)

        # ── Copier le gabarit → destination (préserve tous les styles/médias) ──
        shutil.copy2(str(tpl_path), str(dest))
        doc = Document(str(dest))

        # ── Substitution des signets (bookmarks) ──────────────────────────────
        bookmarks = document.get("bookmarks", {})
        substituted = []
        for bm_name, bm_value in bookmarks.items():
            ok = _replace_bookmark_text(doc, bm_name, str(bm_value))
            if ok:
                substituted.append(bm_name)

        # ── Vider le corps si demandé ─────────────────────────────────────────
        if clear_body:
            body = doc.element.body
            # Conserver uniquement les éléments de mise en page (sectPr)
            to_remove = []
            for child in body:
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag not in ("sectPr",):
                    to_remove.remove if False else to_remove.append(child)
            for elem in to_remove:
                body.remove(elem)

        # ── Titre principal ───────────────────────────────────────────────────
        title = document.get("title", "")
        if title:
            # Chercher un style de titre dans le gabarit
            title_style = "Heading 1"
            for candidate in ("Titre", "Titre 1", "Title", "Heading 1"):
                try:
                    doc.styles[candidate]
                    title_style = candidate
                    break
                except KeyError:
                    continue
            p = doc.add_heading(title, level=0)
            try:
                p.style = title_style
            except Exception:
                pass

        # ── Sections ──────────────────────────────────────────────────────────
        section_count = 0
        for section in document.get("sections", []):
            heading   = section.get("heading", "")
            level     = max(1, min(3, int(section.get("level", 1))))
            custom_style = section.get("style", "")

            if heading:
                # Priorité : style custom > heading Word standard
                if custom_style:
                    p = doc.add_paragraph(heading)
                    _apply_style_safe(p, custom_style, doc) or _apply_style_safe(
                        p, f"Heading {level}", doc
                    )
                else:
                    doc.add_heading(heading, level=level)

            # ── Paragraphes ──────────────────────────────────────────────────
            body_style = custom_style if custom_style and not heading else ""
            # Style de corps par défaut : chercher "Corps texte" ou "Normal"
            if not body_style:
                for candidate in ("Corps texte", "Body Text", "Normal"):
                    try:
                        doc.styles[candidate]
                        body_style = candidate
                        break
                    except KeyError:
                        continue

            if section.get("paragraphs"):
                for para in section["paragraphs"]:
                    if para and str(para).strip():
                        p = doc.add_paragraph(str(para))
                        if body_style:
                            _apply_style_safe(p, body_style, doc)

            elif section.get("content"):
                raw = section["content"]
                parts = (
                    [s.strip() for s in raw.split("\n\n") if s.strip()]
                    if "\n\n" in raw
                    else [s.strip() for s in raw.split("\n") if s.strip()]
                )
                for part in parts:
                    p = doc.add_paragraph(part)
                    if body_style:
                        _apply_style_safe(p, body_style, doc)

            if section.get("intro"):
                p = doc.add_paragraph(str(section["intro"]))
                if body_style:
                    _apply_style_safe(p, body_style, doc)

            # ── Liste à puces ────────────────────────────────────────────────
            if section.get("bullets"):
                blt_style = section.get("bullet_style", "")
                if not blt_style:
                    for candidate in ("Puce", "List Bullet", "Liste à puces"):
                        try:
                            doc.styles[candidate]
                            blt_style = candidate
                            break
                        except KeyError:
                            continue
                if not blt_style:
                    blt_style = "List Bullet"
                for item in section["bullets"]:
                    if item and str(item).strip():
                        p = doc.add_paragraph(str(item))
                        try:
                            p.style = blt_style
                        except Exception:
                            pass

            # ── Tableau ──────────────────────────────────────────────────────
            if section.get("table"):
                tbl_data = section["table"]
                headers  = tbl_data.get("headers", [])
                rows     = tbl_data.get("rows", [])
                tbl_style = section.get("table_style", "")
                if headers:
                    # Chercher un style de tableau dans le gabarit
                    if not tbl_style:
                        for candidate in ("Tableau grille", "Table Grid",
                                          "Light Grid Accent 1", "Light List"):
                            try:
                                doc.styles[candidate]
                                tbl_style = candidate
                                break
                            except KeyError:
                                continue
                    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                    if tbl_style:
                        try:
                            table.style = tbl_style
                        except Exception:
                            table.style = "Table Grid"
                    hdr_cells = table.rows[0].cells
                    for i, h in enumerate(headers):
                        hdr_cells[i].text = str(h)
                    for ri, row in enumerate(rows):
                        row_cells = table.rows[ri + 1].cells
                        for ci, val in enumerate(row[: len(headers)]):
                            row_cells[ci].text = str(val)
                    doc.add_paragraph()

            if section.get("page_break"):
                doc.add_page_break()

            section_count += 1

        doc.save(str(dest))
        extra = {
            "template": str(tpl_path),
            "sections": section_count,
            "bookmarks_substituted": substituted,
        }
        return _ok(dest, extra)

    except Exception as e:
        return _err(f"export_docx_template : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# POWERPOINT — EXPORT DEPUIS GABARIT PPTX
# ═════════════════════════════════════════════════════════════════════════════

@tool(
    name="export_pptx_template",
    description=(
        "Génère une présentation PowerPoint (.pptx) EN APPLIQUANT UN GABARIT ORGANISATIONNEL. "
        "Le thème, les couleurs, les polices, les logos et les layouts du gabarit sont "
        "INTÉGRALEMENT préservés. "
        "C'est l'outil PRIORITAIRE pour toute création de présentation PowerPoint dès "
        "qu'un fichier gabarit (.pptx) est disponible. "
        "Vérifier SYSTÉMATIQUEMENT si un gabarit existe dans ~/Modèles/ avant d'utiliser "
        "export_pptx_json ou export_pptx_outline. "
        "Utiliser cet outil quand : l'utilisateur demande une présentation PowerPoint "
        "(même sans mentionner de gabarit), quand il mentionne un modèle/template/charte "
        "graphique/modèle maison/modèle de l'organisation. "
        "\n\n"
        "FONCTIONNEMENT : "
        "1. Le gabarit est ouvert (thème et layouts chargés). "
        "2. Les slides existants sont supprimés (sauf si keep_example_slides=true). "
        "3. Les nouveaux slides sont créés depuis les layouts du gabarit. "
        "   Spécifier layout_name (par nom) ou layout_index (0-based, défaut 1). "
        "4. Les placeholders sont peuplés avec title, bullets, content ou placeholders. "
        "\n\n"
        "Layouts courants : index 0 = titre/couverture, index 1 = titre+contenu (défaut), "
        "index 5 = titre seul, index 6 = vide. "
        "Les noms exacts sont retournés dans le champ layouts_available de la réponse."
    ),
    parameters={
        "type": "object",
        "properties": {
            "template_path": {
                "type": "string",
                "description": (
                    "Chemin vers le gabarit .pptx de l'organisation. "
                    "Ex : ~/Modèles/presentation_corp.pptx"
                )
            },
            "presentation": {
                "type": "object",
                "description": (
                    "Structure de la présentation à générer. "
                    "Champs de premier niveau : "
                    "title (str, optionnel — pour le slide de titre), "
                    "subtitle (str, optionnel — sous-titre du slide de titre), "
                    "slides (liste d'objets slide). "
                    "Chaque slide accepte : "
                    "title (str), "
                    "layout_name (str — nom exact du layout dans le gabarit, optionnel), "
                    "layout_index (int — index 0-based du layout, défaut 1), "
                    "bullets (liste de str — points principaux), "
                    "content (str — texte libre si pas de bullets), "
                    "notes (str — notes de présentation), "
                    "placeholders (objet idx→texte pour cibler des placeholders précis, "
                    "  ex : {\"0\": \"Titre\", \"1\": \"Corps\", \"2\": \"Date\"})."
                )
            },
            "keep_example_slides": {
                "type": "boolean",
                "description": (
                    "Si false (défaut), les slides d'exemple du gabarit sont supprimés. "
                    "Si true, ils sont conservés et le nouveau contenu est ajouté après."
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination (ex: ~/Documents/présentation.pptx). Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["template_path", "presentation"]
    }
)
def export_pptx_template(
    template_path: str,
    presentation: dict,
    keep_example_slides: bool = False,
    output_path: str = "",
    filename: str = ""
) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        tpl_path = _resolve_template(template_path)
        if not tpl_path.exists():
            return _err(f"export_pptx_template : gabarit introuvable : {tpl_path}")
        if tpl_path.suffix.lower() != ".pptx":
            return _err("export_pptx_template : le gabarit doit être un fichier .pptx")

        title_str = presentation.get("title", "présentation")
        name = filename or (title_str + ".pptx")
        if not name.endswith(".pptx"):
            name += ".pptx"
        dest = _resolve_output(output_path, name)

        # ── Ouvrir le gabarit ─────────────────────────────────────────────────
        prs = Presentation(str(tpl_path))

        # ── Supprimer les slides d'exemple si demandé ─────────────────────────
        if not keep_example_slides:
            xml_slides = prs.slides._sldIdLst
            while len(prs.slides) > 0:
                slide_elem = prs.slides._sldIdLst[0]
                prs.slides._sldIdLst.remove(slide_elem)

        # ── Résolution d'un layout par nom ou index ───────────────────────────
        def _get_layout(name_or_none: str | None, index: int):
            layouts = prs.slide_layouts
            if name_or_none:
                # Cherche par nom exact, puis par correspondance partielle
                for lay in layouts:
                    if lay.name == name_or_none:
                        return lay
                for lay in layouts:
                    if name_or_none.lower() in lay.name.lower():
                        return lay
            try:
                return layouts[index]
            except IndexError:
                return layouts[min(1, len(layouts) - 1)]

        # ── Slide de titre ────────────────────────────────────────────────────
        main_title = presentation.get("title", "")
        subtitle   = presentation.get("subtitle", "")
        if main_title:
            layout = _get_layout(None, 0)
            slide  = prs.slides.add_slide(layout)
            for ph in slide.placeholders:
                idx = ph.placeholder_format.idx
                if idx == 0 and main_title:
                    ph.text = main_title
                elif idx == 1 and subtitle:
                    ph.text = subtitle

        # ── Slides de contenu ─────────────────────────────────────────────────
        slide_count = 0
        for slide_def in presentation.get("slides", []):
            layout_name  = slide_def.get("layout_name")
            layout_index = int(slide_def.get("layout_index", 1))
            layout = _get_layout(layout_name, layout_index)
            slide  = prs.slides.add_slide(layout)

            slide_title   = slide_def.get("title", "")
            bullets       = slide_def.get("bullets", [])
            content_text  = slide_def.get("content", "")
            notes_text    = slide_def.get("notes", "")
            custom_phs    = slide_def.get("placeholders", {})

            # Substitution par idx personnalisé (prioritaire)
            if custom_phs:
                for ph in slide.placeholders:
                    idx_str = str(ph.placeholder_format.idx)
                    if idx_str in custom_phs:
                        try:
                            ph.text = str(custom_phs[idx_str])
                        except Exception:
                            pass
            else:
                # Comportement standard : titre (idx 0) + corps (idx 1 ou 2)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_title

                for ph in slide.placeholders:
                    idx = ph.placeholder_format.idx
                    if idx in (1, 2):
                        tf = ph.text_frame
                        tf.clear()
                        if bullets:
                            for i, b in enumerate(bullets):
                                if i == 0:
                                    tf.paragraphs[0].text  = b
                                    tf.paragraphs[0].level = 0
                                else:
                                    p = tf.add_paragraph()
                                    p.text  = b
                                    p.level = 0
                        elif content_text:
                            tf.paragraphs[0].text = content_text
                        break

            # Notes
            if notes_text:
                try:
                    slide.notes_slide.notes_text_frame.text = notes_text
                except Exception:
                    pass

            slide_count += 1

        prs.save(str(dest))
        layouts_available = [lay.name for lay in prs.slide_layouts]
        extra = {
            "template":          str(tpl_path),
            "slides":            len(prs.slides),
            "layouts_available": layouts_available,
        }
        return _ok(dest, extra)

    except Exception as e:
        return _err(f"export_pptx_template : {e}")
